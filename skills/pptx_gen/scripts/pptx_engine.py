"""
Core PPTX rendering engine using python-pptx.

Provides slide primitives that adapters compose into decks.
"""

import os
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Design system colors
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
SUCCESS = RGBColor(0x28, 0xA7, 0x45)
WARNING = RGBColor(0xFF, 0xC1, 0x07)
DANGER = RGBColor(0xDC, 0x35, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

PRIORITY_COLORS = {
    "Critical": DANGER,
    "High": WARNING,
    "Medium": ACCENT,
    "Low": SUCCESS,
}


def create_presentation() -> Presentation:
    """Create a new blank presentation with standard dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """Add a title slide with navy background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    # Title
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    subtitle: str = "",
) -> None:
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = NAVY
    p.font.bold = True

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_GRAY

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    bf = bullet_box.text_frame
    bf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = bf.paragraphs[0]
        else:
            p = bf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: List[str],
    right_title: str,
    right_bullets: List[str],
) -> None:
    """Add a two-column layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = NAVY
    p.font.bold = True

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
    lf = left_box.text_frame
    lf.word_wrap = True
    p = lf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    for bullet in left_bullets:
        p = lf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY

    # Right column
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.3), Inches(5.8), Inches(5.5))
    rf = right_box.text_frame
    rf.word_wrap = True
    p = rf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    for bullet in right_bullets:
        p = rf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: List[str],
    rows: List[List[str]],
) -> None:
    """Add a data table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = NAVY
    p.font.bold = True

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.5),
        Inches(12), Inches(5),
    )
    table = table_shape.table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)


def save_presentation(prs: Presentation, output_path: str) -> str:
    """Save the presentation to disk."""
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path
